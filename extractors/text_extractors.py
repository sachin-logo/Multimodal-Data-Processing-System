import os
from PyPDF2 import PdfReader
import docx
from pptx import Presentation

# PDF extraction
def extract_pdf_text(path):
    reader = PdfReader(path)
    return '\n'.join([page.extract_text() for page in reader.pages])

# DOCX extraction
def extract_docx_text(path):
    doc = docx.Document(path)
    return '\n'.join([para.text for para in doc.paragraphs])

# PPTX extraction
def extract_pptx_text(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return '\n'.join(text)

# TXT/MD extraction
def extract_txt_md(path):
    # Try common encodings with graceful fallback
    encodings_to_try = ['utf-8', 'utf-8-sig', 'cp1252', 'latin-1']
    for enc in encodings_to_try:
        try:
            with open(path, 'r', encoding=enc, errors='strict') as f:
                return f.read()
        except UnicodeDecodeError:
            continue
        except Exception:
            continue
    # Last resort: read as bytes and decode replacing errors
    try:
        with open(path, 'rb') as f:
            data = f.read()
        return data.decode('utf-8', errors='replace')
    except Exception:
        # Fallback minimal message to avoid crashing the pipeline
        return ""
